import os
import ssl
import time
import tempfile
import urllib3
import concurrent.futures

# Disable SSL verification for corporate proxy/firewall
os.environ["PYTHONHTTPSVERIFY"] = "0"
os.environ["CURL_CA_BUNDLE"] = ""
os.environ["SSL_CERT_FILE"] = ""
os.environ["REQUESTS_CA_BUNDLE"] = ""
ssl._create_default_https_context = ssl._create_unverified_context
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

import io
import json
import re
import streamlit as st
import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageOps
from google import genai
from google.genai import types

# Optional PDF parsing
try:
    import pypdf
except ImportError:
    pypdf = None

# Optional Text-to-Speech Audio Generation
try:
    from gtts import gTTS
except ImportError:
    gTTS = None

# Optional MP4 Video Generation Engine (Supporting MoviePy v1 and v2)
MOVIEPY_AVAILABLE = False
MOVIEPY_ERROR = ""
try:
    try:
        from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
        MOVIEPY_AVAILABLE = True
    except Exception:
        from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
        MOVIEPY_AVAILABLE = True
except Exception as e:
    MOVIEPY_AVAILABLE = False
    MOVIEPY_ERROR = str(e)

# ---------------------------------------------------------
# Primary Model Cascade (Gemini 3.x Flash Family)
# ---------------------------------------------------------
MODEL_CASCADE = [
    "gemini-3.7-flash",       # ⭐ Best default (General AI, reasoning, coding)
    "gemini-3.6-flash",       # ⭐ Fast general-purpose tasks
    "gemini-3.5-flash",       # General-purpose AI
    "gemini-3.5-flash-lite",  # 💰 Cost-focused, high volume
    "gemini-3.1-flash-lite"   # Lightweight fallback
]

def generate_resilient_content(client_inst, contents, system_prompt="", response_mime_type=None, max_retries=2):
    """
    Executes generation with auto-retry and cascading fallback across
    the requested Gemini 3.x Flash model series.
    """
    config_args = {"temperature": 0.2}
    if system_prompt:
        config_args["system_instruction"] = system_prompt
    if response_mime_type:
        config_args["response_mime_type"] = response_mime_type
    
    config = types.GenerateContentConfig(**config_args)
    last_err = None

    for model_name in MODEL_CASCADE:
        for attempt in range(max_retries):
            try:
                resp = client_inst.models.generate_content(
                    model=model_name,
                    contents=contents,
                    config=config
                )
                if resp and resp.text:
                    return resp.text, model_name
            except Exception as e:
                err_str = str(e).lower()
                last_err = e
                if any(code in err_str for code in ["503", "unavailable", "demand", "429", "quota", "resource_exhausted"]):
                    time.sleep(1.5 * (attempt + 1))
                    continue
                elif "404" in err_str or "not found" in err_str:
                    break
                else:
                    time.sleep(1.0)
                    
    raise last_err

# ---------------------------------------------------------
# Page Configuration & Visual Theme
# ---------------------------------------------------------
st.set_page_config(
    page_title="ReqAssist - BA/PM Requirement Engine",
    page_icon="🚀",
    layout="wide"
)

# Custom CSS
st.markdown("""
<style>
    .block-container {
        padding-top: 3.5rem !important;
        padding-bottom: 4rem !important;
    }
    
    section[data-testid="stSidebar"] h1 {
        font-size: 1.25rem !important;
        white-space: nowrap !important;
        overflow: hidden !important;
        text-overflow: ellipsis !important;
        line-height: 1.3 !important;
        padding-bottom: 0.4rem !important;
    }

    div.stButton > button {
        width: 100% !important;
        height: 38px !important;
        min-height: 38px !important;
        max-height: 38px !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        border-radius: 6px !important;
        font-size: 14px !important;
        font-weight: 500 !important;
        padding: 0 12px !important;
        margin: 0 !important;
        line-height: 1 !important;
        white-space: nowrap !important;
        text-align: center !important;
        background-color: #ffffff !important;
        color: #1e293b !important;
        border: 1px solid #cbd5e1 !important;
        box-shadow: 0 1px 2px rgba(0, 0, 0, 0.05) !important;
        transition: all 0.2s ease-in-out !important;
    }

    div.stButton > button p,
    div.stButton > button div,
    div.stButton > button span {
        font-size: 14px !important;
        font-weight: 500 !important;
        margin: 0 !important;
        padding: 0 !important;
        line-height: 1 !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        color: inherit !important;
    }

    div.stButton > button:hover {
        background-color: #f1f5f9 !important;
        border-color: #94a3b8 !important;
        color: #0f172a !important;
    }

    div.stButton:not(.st-key-main_generate_btn) > button[kind="primary"] {
        background-color: #2563eb !important;
        color: #ffffff !important;
        border: 1.5px solid #1d4ed8 !important;
        box-shadow: 0 2px 6px rgba(37, 99, 235, 0.30) !important;
    }

    div.st-key-btn_lang_en button::before {
        content: "" !important;
        display: inline-block !important;
        width: 18px !important;
        height: 12px !important;
        margin-right: 6px !important;
        background-image: url('https://flagcdn.com/w40/gb.png') !important;
        background-size: cover !important;
        background-position: center !important;
        border-radius: 2px !important;
    }

    div.st-key-btn_lang_it button::before {
        content: "" !important;
        display: inline-block !important;
        width: 18px !important;
        height: 12px !important;
        margin-right: 6px !important;
        background-image: url('https://flagcdn.com/w40/it.png') !important;
        background-size: cover !important;
        background-position: center !important;
        border-radius: 2px !important;
    }

    .st-key-main_generate_btn button,
    .st-key-main_generate_btn button[kind="primary"],
    div[data-testid="stButton"].st-key-main_generate_btn > button,
    div[data-testid="stButton"].st-key-main_generate_btn > button[kind="primary"] {
        background-color: #051330 !important;
        color: #ffffff !important;
        border: 1px solid #051330 !important;
        font-size: 22px !important;
        font-weight: 700 !important;
        height: 64px !important;
        min-height: 64px !important;
        max-height: 64px !important;
        padding: 0 40px !important;
        border-radius: 10px !important;
        box-shadow: 0 6px 20px rgba(5, 19, 48, 0.45) !important;
        transition: all 0.25s cubic-bezier(0.4, 0, 0.2, 1) !important;
    }

    .st-key-main_generate_btn button p,
    .st-key-main_generate_btn button div,
    .st-key-main_generate_btn button span {
        color: #ffffff !important;
        font-weight: 700 !important;
        font-size: 22px !important;
        line-height: 1.2 !important;
    }

    .st-key-main_generate_btn button:hover,
    .st-key-main_generate_btn button:active,
    .st-key-main_generate_btn button:focus:active,
    div[data-testid="stButton"].st-key-main_generate_btn > button:hover,
    div[data-testid="stButton"].st-key-main_generate_btn > button:active {
        background-color: #f56642 !important;
        border-color: #f56642 !important;
        color: #ffffff !important;
        transform: translateY(-2px);
        box-shadow: 0 10px 28px rgba(245, 102, 66, 0.45) !important;
    }

    .plain-upload-status {
        color: #15803d !important;
        font-size: 14px !important;
        font-weight: 600 !important;
        margin-top: 6px !important;
        background: transparent !important;
        border: none !important;
        padding: 0 !important;
    }
    .plain-upload-filename {
        color: #166534 !important;
        font-weight: 700 !important;
    }
</style>
""", unsafe_allow_html=True)

# ---------------------------------------------------------
# Centered Orange Circular Progress (% in center)
# ---------------------------------------------------------
def render_circular_progress(placeholder, percent: int, label: str = "Generating..."):
    """Renders a centered orange circular progress ring starting from 1% with percentage in center."""
    percent = max(1, min(100, int(percent)))
    circumference = 314.16
    offset = circumference - (circumference * percent / 100)
    
    html_code = f"""
    <div style="display: flex; flex-direction: column; align-items: center; justify-content: center; width: 100%; margin: 1.5rem auto 2rem auto;">
        <div style="position: relative; width: 130px; height: 130px; display: flex; align-items: center; justify-content: center;">
            <svg style="transform: rotate(-90deg); width: 130px; height: 130px;">
                <circle cx="65" cy="65" r="50" stroke="#e2e8f0" stroke-width="8" fill="transparent" />
                <circle cx="65" cy="65" r="50" stroke="#f56642" stroke-width="8" fill="transparent"
                    stroke-dasharray="{circumference}"
                    stroke-dashoffset="{offset}"
                    stroke-linecap="round"
                    style="transition: stroke-dashoffset 0.1s linear;" />
            </svg>
            <div style="position: absolute; text-align: center;">
                <span style="font-size: 24px; font-weight: 800; color: #051330; font-family: sans-serif;">{percent}%</span>
            </div>
        </div>
        <p style="margin-top: 12px; font-size: 16px; font-weight: 600; color: #334155; font-family: sans-serif; text-align: center;">
            {label}
        </p>
    </div>
    """
    placeholder.markdown(html_code, unsafe_allow_html=True)

# ---------------------------------------------------------
# UI Dictionary (English & Italiano)
# ---------------------------------------------------------
T = {
    "English": {
        "sidebar_title": "🚀 ReqAssist Settings",
        "lang_label": "Language / Lingua:",
        "name_label": "Your Name:",
        "role_label": "Select Your Role:",
        "roles": ["BA / PM / PO (Full Authoring)", "Developer / Tester / Viewer"],
        "api_key_warn": "⚠️ Gemini API Key missing. Please configure GEMINI_API_KEY in Streamlit Cloud Secrets.",
        "welcome": "Let's turn your raw requirements into high-impact deliverables.",
        "step1_title": "📂 STEP 1: Upload Source Documentation & Assets",
        "raw_title": "📄 Functional Notes RAW",
        "excel_title": "📊 Field Validations",
        "figma_title": "🎨 Figma UI Screens (.png, .jpg)",
        "supported_formats_help": "Supported: .docx, .txt, .pdf, .xlsx, .csv, .xls",
        "step2_title": "⚡ STEP 2: Select & Generate Deliverable",
        "generate_btn": "✨ Generate Deliverable",
        "tab_output": "📄 Output",
        "tab_download": "📥 Download Formatted Files",
        "docx_btn": "📄 Download Word (.docx)",
        "xlsx_btn": "📊 Download Test Cases (.xlsx)",
        "pptx_btn": "📊 Download Presentation (.pptx)",
        "md_btn": "📝 Download Markdown (.md)",
        "audio_btn": "🎙️ Download Demo Voiceover (.mp3)",
        "mp4_btn": "🎬 Download Demo Video (.mp4)",
        "video_player_title": "🎬 Generated MP4 Demo Walkthrough (Watch & Download)",
        "audio_player_title": "🎧 Generated Voiceover Audio Track (Listen & Download)",
        "viewer_err": "🚫 Access Restricted: These options are restricted to BA & PM roles. Please contact your project PM or PO.",
        "popup_title": "⚠️ Missing Required Documents",
        "popup_msg": "ReqAssist strictly requires the following document(s) before generating this artifact. Please upload them in Step 1:",
        "popup_btn": "OK, Got It",
        "quiz_header": "🧠 Interactive Requirements & Knowledge Check",
        "quiz_submit_btn": "🎯 Submit Quiz for Grading",
        "quiz_retake_btn": "🔄 Retake Quiz",
        "quiz_score_title": "Quiz Assessment Results",
        "quiz_passed": "🎉 Excellent work! You passed the requirements check (≥ 75%)!",
        "quiz_failed": "📚 Score is below 75%. Please review the detailed explanations below:",
        "quiz_focus_title": "🎯 Recommended Areas to Focus On:",
        "quiz_focus_intro": "Pay close attention to these requirement topics where questions were missed:",
        "quiz_your_ans": "Your Answer:",
        "quiz_correct_ans": "Correct Answer:",
        "quiz_explanation": "Explanation & Requirement Reference:",
        "options": [
            "📝 Acceptance Criteria (AC)",
            "🧪 Test Cases",
            "📑 Detailed Functional Analysis",
            "🎥 Demo Video",
            "📊 Presentation / PPT (5 to 10 slides)",
            "❓ FAQs",
            "🧠 Quiz"
        ]
    },
    "Italiano": {
        "sidebar_title": "🚀 Impostazioni ReqAssist",
        "lang_label": "Language / Lingua:",
        "name_label": "Il Tuo Nome:",
        "role_label": "Seleziona il Tuo Ruolo:",
        "roles": ["BA / PM / PO (Autore Completo)", "Sviluppatore / Tester / Viewer"],
        "api_key_warn": "⚠️ Chiave API Gemini mancante. Configura GEMINI_API_KEY nei Secrets di Streamlit Cloud.",
        "welcome": "Trasformiamo i tuoi requisiti in deliverable di alto impatto.",
        "step1_title": "📂 PASSAGGIO 1: Carica la Documentazione Sorgente & Asset",
        "raw_title": "📄 Analisi Funzionale RAW",
        "excel_title": "📊 Validazioni Campi",
        "figma_title": "🎨 Schermate UI Figma (.png, .jpg)",
        "supported_formats_help": "Formati supportati: .docx, .txt, .pdf, .xlsx, .csv, .xls",
        "step2_title": "⚡ PASSAGGIO 2: Seleziona & Genera Deliverable",
        "generate_btn": "✨ Genera Deliverable",
        "tab_output": "📄 Output",
        "tab_download": "📥 Scarica File Formattati",
        "docx_btn": "📄 Scarica Documento Word (.docx)",
        "xlsx_btn": "📊 Scarica Test Case in Excel (.xlsx)",
        "pptx_btn": "📊 Scarica Presentazione (.pptx)",
        "md_btn": "📝 Scarica Markdown (.md)",
        "audio_btn": "🎙️ Scarica Traccia Audio Demo (.mp3)",
        "mp4_btn": "🎬 Scarica Video Demo (.mp4)",
        "video_player_title": "🎬 Video Demo MP4 Generato (Guarda & Scarica)",
        "audio_player_title": "🎧 Traccia Audio Voiceover Generata (Ascolta & Scarica)",
        "viewer_err": "🚫 Accesso Limitato: Queste opzioni sono riservate a BA e PM. Contatta il PM o PO del progetto.",
        "popup_title": "⚠️ Documenti Obbligatori Mancanti",
        "popup_msg": "ReqAssist richiede obbligatoriamente i seguenti documenti prima di procedere. Caricali nel Passaggio 1:",
        "popup_btn": "OK, Ho Capito",
        "quiz_header": "🧠 Quiz Interattivo di Verifica Requisiti",
        "quiz_submit_btn": "🎯 Invia Quiz per la Valutazione",
        "quiz_retake_btn": "🔄 Ripeti il Quiz",
        "quiz_score_title": "Risultati della Valutazione del Quiz",
        "quiz_passed": "🎉 Ottimo lavoro! Hai superato la verifica dei requisiti (≥ 75%)!",
        "quiz_failed": "📚 Punteggio inferiore al 75%. Consulta le spiegazioni dettagliate di seguito:",
        "quiz_focus_title": "🎯 Aree di Focus Consigliate:",
        "quiz_focus_intro": "Presta particolare attenzione a questi requisiti in cui le risposte non erano corrette:",
        "quiz_your_ans": "La Tua Risposta:",
        "quiz_correct_ans": "Risposta Corretta:",
        "quiz_explanation": "Spiegazione e Riferimento ai Requisiti:",
        "options": [
            "📝 Acceptance Criteria (AC)",
            "🧪 Test Cases (Casi di Test)",
            "📑 Analisi Funzionale Dettagliata",
            "🎥 Video Demo",
            "📊 Presentazione / PPT (5-10 slide)",
            "❓ FAQ (Domande Frequenti)",
            "🧠 Quiz"
        ]
    }
}

# ---------------------------------------------------------
# Sidebar Setup: Language & Role Selection
# ---------------------------------------------------------
if "selected_lang" not in st.session_state:
    st.session_state.selected_lang = "English"

st.sidebar.markdown("**Language / Lingua:**")

col_en, col_it = st.sidebar.columns(2)
with col_en:
    uk_selected = (st.session_state.selected_lang == "English")
    if st.button("English", key="btn_lang_en", type="primary" if uk_selected else "secondary", use_container_width=True):
        st.session_state.selected_lang = "English"
        st.rerun()

with col_it:
    it_selected = (st.session_state.selected_lang == "Italiano")
    if st.button("Italiano", key="btn_lang_it", type="primary" if it_selected else "secondary", use_container_width=True):
        st.session_state.selected_lang = "Italiano"
        st.rerun()

lang_key = st.session_state.selected_lang
ui = T[lang_key]

st.sidebar.title(ui["sidebar_title"])
user_name = st.sidebar.text_input(ui["name_label"], value="Mayank", key="user_name_input")
user_role = st.sidebar.selectbox(ui["role_label"], ui["roles"], key="user_role_select")

api_key = st.secrets.get("GEMINI_API_KEY", os.environ.get("GEMINI_API_KEY", ""))

# ---------------------------------------------------------
# Header & Welcome
# ---------------------------------------------------------
st.title("🚀 ReqAssist")
st.caption(f"**{'Benvenuto' if lang_key == 'Italiano' else 'Welcome'}, {user_name}!** {ui['welcome']}")

if not api_key:
    st.error(ui["api_key_warn"])
    st.stop()

client = genai.Client(
    api_key=api_key,
    http_options=types.HttpOptions(
        client_args={"verify": False}
    )
)

# ---------------------------------------------------------
# Modal Alert Popup Component
# ---------------------------------------------------------
@st.dialog(ui["popup_title"])
def show_missing_data_popup(missing_list):
    st.error(ui["popup_msg"])
    for item in missing_list:
        st.markdown(f"❌ **{item}**")
    st.write("")
    if st.button(ui["popup_btn"], type="primary", use_container_width=True):
        st.rerun()

# ---------------------------------------------------------
# Universal Multi-Format File Reader
# ---------------------------------------------------------
def parse_uploaded_file(uploaded_file) -> str:
    """Extracts clean text or tabular representations from docx, pdf, xlsx, xls, csv, txt, md."""
    if not uploaded_file:
        return ""
    name = uploaded_file.name.lower()
    try:
        if name.endswith(".docx"):
            doc = Document(uploaded_file)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        
        elif name.endswith(".pdf"):
            if pypdf is None:
                return "PDF parser (pypdf) is not installed."
            reader = pypdf.PdfReader(uploaded_file)
            pages_text = []
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    pages_text.append(extracted)
            return "\n".join(pages_text)
        
        elif name.endswith(".xlsx") or name.endswith(".xls"):
            df = pd.read_excel(uploaded_file)
            return df.to_string(index=False)
            
        elif name.endswith(".csv"):
            df = pd.read_csv(uploaded_file)
            return df.to_string(index=False)
            
        else:
            return uploaded_file.read().decode("utf-8", errors="ignore")
            
    except Exception as e:
        return f"Error reading {uploaded_file.name}: {str(e)}"

# ---------------------------------------------------------
# Helper Functions: Word DOCX Exporter
# ---------------------------------------------------------
def create_docx(title: str, content: str) -> io.BytesIO:
    """Creates a well-formatted MS Word document (.docx) from markdown text."""
    doc = Document()
    doc.add_heading(title, 0)
    for line in content.split("\n"):
        if line.startswith("## "):
            doc.add_heading(line.replace("## ", ""), level=1)
        elif line.startswith("### "):
            doc.add_heading(line.replace("### ", ""), level=2)
        elif line.startswith("* ") or line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif line.strip():
            doc.add_paragraph(line)
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio

# ---------------------------------------------------------
# Context-Aware AI Diagram Engine for PowerPoint & Video
# ---------------------------------------------------------
def generate_contextual_ai_slide_diagram(slide_title: str, bullets: list) -> io.BytesIO:
    """
    Generates a tailored visual diagram based on the SPECIFIC content of the slide.
    Categorizes context (API, DB, Auth, Workflow, UI) and creates matching graphics.
    """
    combined_text = f"{slide_title} " + " ".join(bullets).lower()
    
    img = Image.new("RGB", (1280, 720), color=(5, 19, 48))
    draw = ImageDraw.Draw(img)
    
    draw.rounded_rectangle([25, 25, 1255, 695], radius=16, outline=(37, 99, 235), width=3)
    draw.rounded_rectangle([45, 45, 1235, 135], radius=12, fill=(15, 30, 65), outline=(59, 130, 246), width=1)
    draw.text((70, 75), f"CONTENT ARCHITECTURE • {slide_title.upper()[:40]}", fill=(255, 255, 255))
    
    b1_text = bullets[0] if len(bullets) > 0 else "System Initialization"
    b2_text = bullets[1] if len(bullets) > 1 else "Business Rule Processing"
    b3_text = bullets[2] if len(bullets) > 2 else "Final Output & State Commit"

    if any(k in combined_text for k in ["api", "endpoint", "payload", "http", "contract", "json"]):
        draw.rounded_rectangle([65, 175, 415, 480], radius=12, fill=(10, 25, 55), outline=(245, 102, 66), width=2)
        draw.text((85, 205), "REQUEST / PAYLOAD", fill=(245, 102, 66))
        draw.text((85, 250), f"• Ingest: {b1_text[:30]}\n• Headers & Tokens\n• JSON Schema Guard", fill=(226, 232, 240))
        
        draw.rounded_rectangle([455, 175, 825, 480], radius=12, fill=(10, 25, 55), outline=(34, 197, 94), width=2)
        draw.text((475, 205), "INTEGRATION SERVICE", fill=(34, 197, 94))
        draw.text((475, 250), f"• Logic: {b2_text[:30]}\n• Data Validation Rules\n• Subsystem Routing", fill=(226, 232, 240))
        
        draw.rounded_rectangle([865, 175, 1215, 480], radius=12, fill=(10, 25, 55), outline=(168, 85, 247), width=2)
        draw.text((885, 205), "RESPONSE CODES", fill=(192, 132, 252))
        draw.text((885, 250), f"• Result: {b3_text[:30]}\n• 200 OK / 400 Bad Request\n• Audit Event Telemetry", fill=(226, 232, 240))

    elif any(k in combined_text for k in ["database", "db", "table", "schema", "storage", "sql", "record"]):
        draw.rounded_rectangle([65, 175, 415, 480], radius=12, fill=(10, 25, 55), outline=(59, 130, 246), width=2)
        draw.text((85, 205), "ENTITY INGESTION", fill=(59, 130, 246))
        draw.text((85, 250), f"• Input: {b1_text[:30]}\n• Type Cast & Constraints\n• Foreign Key Checks", fill=(226, 232, 240))
        
        draw.rounded_rectangle([455, 175, 825, 480], radius=12, fill=(10, 25, 55), outline=(245, 102, 66), width=2)
        draw.text((475, 205), "PERSISTENCE & COMMIT", fill=(245, 102, 66))
        draw.text((475, 250), f"• Transaction: {b2_text[:30]}\n• ACID Compliant Writes\n• Rollback on Exception", fill=(226, 232, 240))
        
        draw.rounded_rectangle([865, 175, 1215, 480], radius=12, fill=(10, 25, 55), outline=(34, 197, 94), width=2)
        draw.text((885, 205), "DATA SYNC & VIEWS", fill=(34, 197, 94))
        draw.text((885, 250), f"• Sync: {b3_text[:30]}\n• Read-Model Updates\n• Cache Invalidation", fill=(226, 232, 240))

    elif any(k in combined_text for k in ["security", "auth", "token", "permission", "role", "login"]):
        draw.rounded_rectangle([65, 175, 415, 480], radius=12, fill=(10, 25, 55), outline=(239, 68, 68), width=2)
        draw.text((85, 205), "AUTH VERIFICATION", fill=(239, 68, 68))
        draw.text((85, 250), f"• Check: {b1_text[:30]}\n• Bearer JWT Extraction\n• Expiry & Signature", fill=(226, 232, 240))
        
        draw.rounded_rectangle([455, 175, 825, 480], radius=12, fill=(10, 25, 55), outline=(245, 102, 66), width=2)
        draw.text((475, 205), "ROLE MATRIX (RBAC)", fill=(245, 102, 66))
        draw.text((475, 250), f"• Scope: {b2_text[:30]}\n• Admin / User / Viewer\n• Access Enforcement", fill=(226, 232, 240))
        
        draw.rounded_rectangle([865, 175, 1215, 480], radius=12, fill=(10, 25, 55), outline=(34, 197, 94), width=2)
        draw.text((885, 205), "SESSION STATE", fill=(34, 197, 94))
        draw.text((885, 250), f"• Granted: {b3_text[:30]}\n• Secure Context Active\n• Audit Logging", fill=(226, 232, 240))

    else:
        draw.rounded_rectangle([65, 175, 415, 480], radius=12, fill=(10, 25, 55), outline=(59, 130, 246), width=2)
        draw.text((85, 205), "STAGE 1 • PRECONDITIONS", fill=(59, 130, 246))
        draw.text((85, 250), f"• Step: {b1_text[:30]}\n• UI Initial State\n• Mandatory Inputs", fill=(226, 232, 240))
        
        draw.rounded_rectangle([455, 175, 825, 480], radius=12, fill=(10, 25, 55), outline=(34, 197, 94), width=2)
        draw.text((475, 205), "STAGE 2 • LOGIC EXECUTION", fill=(34, 197, 94))
        draw.text((475, 250), f"• Rule: {b2_text[:30]}\n• Data Transformation\n• State Machine Update", fill=(226, 232, 240))
        
        draw.rounded_rectangle([865, 175, 1215, 480], radius=12, fill=(10, 25, 55), outline=(168, 85, 247), width=2)
        draw.text((885, 205), "STAGE 3 • OUTCOME / UI", fill=(192, 132, 252))
        draw.text((885, 250), f"• Output: {b3_text[:30]}\n• Feedback & Navigation\n• Success Confirmation", fill=(226, 232, 240))
    
    draw.rectangle([415, 320, 455, 326], fill=(245, 102, 66))
    draw.rectangle([825, 320, 865, 326], fill=(34, 197, 94))
    
    draw.rounded_rectangle([900, 625, 1235, 675], radius=8, fill=(15, 30, 65), outline=(245, 102, 66), width=2)
    draw.text((925, 642), "ReqAssist • AI Generated Visual", fill=(245, 102, 66))
    
    bio = io.BytesIO()
    img.save(bio, format="PNG")
    bio.seek(0)
    return bio

def add_watermark_to_figma_screen(uploaded_file) -> io.BytesIO:
    """Overlays the official ReqAssist watermark badge onto uploaded Figma screens."""
    try:
        uploaded_file.seek(0)
        img = Image.open(uploaded_file).convert("RGBA")
        
        badge_w, badge_h = 260, 45
        badge = Image.new("RGBA", (badge_w, badge_h), (5, 19, 48, 230))
        b_draw = ImageDraw.Draw(badge)
        b_draw.rectangle([0, 0, badge_w-1, badge_h-1], outline=(245, 102, 66, 255), width=2)
        b_draw.text((15, 12), "ReqAssist • Figma Visual", fill=(255, 255, 255, 255))
        
        img.paste(badge, (max(10, img.width - badge_w - 20), max(10, img.height - badge_h - 20)), badge)
        
        bio = io.BytesIO()
        img.save(bio, format="PNG")
        bio.seek(0)
        return bio
    except Exception:
        uploaded_file.seek(0)
        return uploaded_file

def create_pptx_with_images(slides_json: list, figma_files: list = None) -> io.BytesIO:
    """
    Generates a 16:9 Presentation:
    - Dark Blue Background (#051330)
    - White Font (#FFFFFF)
    - Integrates AI Content-Driven Visual Diagrams AND Figma screens with ReqAssist watermark
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    DARK_BLUE = RGBColor(5, 19, 48)
    WHITE = RGBColor(255, 255, 255)
    SOFT_WHITE = RGBColor(241, 245, 249)
    MUTED_GREY = RGBColor(148, 163, 184)
    
    total_slides = len(slides_json)
    visual_indices = set()
    num_visual_slides = max(int(total_slides * 0.5), 2)
    step = max(total_slides // num_visual_slides, 1)
    for i in range(0, total_slides, step):
        visual_indices.add(i)
        if len(visual_indices) >= num_visual_slides:
            break
            
    figma_idx = 0
    prepared_figma = [add_watermark_to_figma_screen(f) for f in (figma_files or [])]
    use_ai_diagram_toggle = True
    
    for slide_idx, slide_data in enumerate(slides_json):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BLUE
        bg.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.85))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_data.get("title", f"Slide {slide_idx + 1}")
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        
        is_visual_slide = (slide_idx in visual_indices)
        text_width = Inches(6.2) if is_visual_slide else Inches(11.5)
        
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), text_width, Inches(5.0))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        bullets = slide_data.get("bullets", [])
        for b_idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = SOFT_WHITE
            p.space_after = Pt(10)
            
        if is_visual_slide:
            try:
                if prepared_figma and not use_ai_diagram_toggle:
                    img_stream = prepared_figma[figma_idx % len(prepared_figma)]
                    figma_idx += 1
                    use_ai_diagram_toggle = True
                else:
                    img_stream = generate_contextual_ai_slide_diagram(slide_data.get("title", "Architecture"), bullets)
                    use_ai_diagram_toggle = False if prepared_figma else True
                
                img_stream.seek(0)
                slide.shapes.add_picture(img_stream, Inches(7.3), Inches(1.5), width=Inches(5.2))
            except Exception:
                pass
                
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
        tf_footer = footer_box.text_frame
        p_footer = tf_footer.paragraphs[0]
        p_footer.text = f"ReqAssist 🚀 • Requirements Architecture | Slide {slide_idx + 1} of {total_slides}"
        p_footer.font.size = Pt(11)
        p_footer.font.color.rgb = MUTED_GREY
        
        if "notes" in slide_data and slide_data["notes"]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

def parse_markdown_tables_to_excel(markdown_text: str) -> io.BytesIO:
    """Parses multiple markdown tables into separate Excel sheets (Functional & Technical)."""
    tables = []
    current_table = []
    in_table = False
    
    for line in markdown_text.split("\n"):
        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            current_table.append(stripped)
            in_table = True
        else:
            if in_table and len(current_table) >= 3:
                tables.append(current_table)
            current_table = []
            in_table = False
    if in_table and len(current_table) >= 3:
        tables.append(current_table)
        
    if not tables:
        return None
        
    bio = io.BytesIO()
    with pd.ExcelWriter(bio, engine="openpyxl") as writer:
        sheet_labels = ["Functional Test Cases", "Technical Test Cases"]
        for idx, tbl_lines in enumerate(tables):
            headers = [c.strip() for c in tbl_lines[0].split("|")[1:-1]]
            data = []
            for row_line in tbl_lines[2:]:
                row = [c.strip() for c in row_line.split("|")[1:-1]]
                if len(row) == len(headers):
                    data.append(row)
            if data:
                df = pd.DataFrame(data, columns=headers)
                sname = sheet_labels[idx] if idx < len(sheet_labels) else f"Test Cases Part {idx+1}"
                df.to_excel(writer, index=False, sheet_name=sname[:31])
    bio.seek(0)
    return bio

def extract_voiceover_script(markdown_text: str) -> str:
    """Extracts spoken narration text from the Demo Video Markdown storyboard."""
    lines = markdown_text.split("\n")
    voiceover_lines = []
    for line in lines:
        if "|" in line:
            parts = [p.strip() for p in line.split("|") if p.strip()]
            if len(parts) >= 3 and not any(h in parts[-1].lower() for h in ["voiceover", "script", "---", "testo"]):
                voiceover_lines.append(parts[-1])
        elif not line.startswith("#") and len(line.strip()) > 20:
            voiceover_lines.append(line.strip())
            
    narration = " ".join(voiceover_lines)
    narration = re.sub(r'[*_#`]', '', narration)
    return narration.strip() if len(narration.strip()) > 10 else "Welcome to the product feature demo walkthrough generated by ReqAssist."

def extract_voiceover_and_synthesize_audio(markdown_text: str, lang: str = "en") -> io.BytesIO:
    """Generates MP3 voiceover audio track from demo storyboard."""
    if gTTS is None:
        return None
    try:
        narration = extract_voiceover_script(markdown_text)
        tts = gTTS(text=narration[:4000], lang="it" if lang == "Italiano" else "en", slow=False)
        audio_bio = io.BytesIO()
        tts.write_to_fp(audio_bio)
        audio_bio.seek(0)
        return audio_bio
    except Exception:
        return None

def normalize_image_to_16_9(img_stream_or_file) -> str:
    """Converts any image into a uniform 1280x720 16:9 canvas with dark blue padding."""
    try:
        img_stream_or_file.seek(0)
        im = Image.open(img_stream_or_file).convert("RGB")
    except Exception:
        im = Image.new("RGB", (1280, 720), color=(5, 19, 48))
        
    canvas = Image.new("RGB", (1280, 720), (5, 19, 48))
    im.thumbnail((1200, 660), Image.Resampling.LANCZOS)
    x = (1280 - im.width) // 2
    y = (720 - im.height) // 2
    canvas.paste(im, (x, y))
    
    draw = ImageDraw.Draw(canvas)
    draw.rounded_rectangle([1000, 660, 1250, 705], radius=6, fill=(15, 30, 65), outline=(245, 102, 66), width=2)
    draw.text((1020, 675), "ReqAssist 🚀 Demo", fill=(255, 255, 255))
    
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    canvas.save(tmp.name, format="PNG")
    tmp.close()
    return tmp.name

def set_clip_duration(clip, duration):
    """MoviePy v1 vs v2 compatibility helper for duration."""
    return clip.with_duration(duration) if hasattr(clip, "with_duration") else clip.set_duration(duration)

def set_clip_audio(video_clip, audio_clip):
    """MoviePy v1 vs v2 compatibility helper for audio."""
    return video_clip.with_audio(audio_clip) if hasattr(video_clip, "with_audio") else video_clip.set_audio(audio_clip)

def generate_mp4_video_robust(markdown_text: str, figma_files: list, lang: str = "en"):
    """
    Directly compiles an MP4 video combining content-aware AI scene cards 
    with uploaded Figma screens, synced to the AI voiceover.
    """
    if not MOVIEPY_AVAILABLE:
        return None, f"MoviePy engine not initialized. ({MOVIEPY_ERROR})"
    if gTTS is None:
        return None, "gTTS audio synthesizer is not installed."
        
    temp_audio_path = None
    temp_video_path = None
    temp_img_paths = []
    
    try:
        narration = extract_voiceover_script(markdown_text)
        temp_audio = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
        temp_audio_path = temp_audio.name
        temp_audio.close()
        
        tts = gTTS(text=narration[:4000], lang="it" if lang == "Italiano" else "en", slow=False)
        tts.save(temp_audio_path)
        
        audio_clip = AudioFileClip(temp_audio_path)
        total_duration = max(float(audio_clip.duration), 5.0)

        # 1. AI Overview Card
        ai_card_1 = generate_contextual_ai_slide_diagram("Demo Architecture Overview", ["Functional Workflow Walkthrough", "UI Interaction Flow", "Data & Contract State"])
        temp_img_paths.append(normalize_image_to_16_9(ai_card_1))

        # 2. Figma UI Screens
        if figma_files and len(figma_files) > 0:
            for f in figma_files:
                p = normalize_image_to_16_9(f)
                temp_img_paths.append(p)

        # 3. AI Execution & Logic Card
        ai_card_2 = generate_contextual_ai_slide_diagram("Validation & Backend Processing", ["Validation Rules & Error Checks", "State Transitions & DB Commit", "Success Response Code"])
        temp_img_paths.append(normalize_image_to_16_9(ai_card_2))

        duration_per_slide = total_duration / len(temp_img_paths)
        clips = [set_clip_duration(ImageClip(p), duration_per_slide) for p in temp_img_paths]
        
        video_clip = concatenate_videoclips(clips, method="compose")
        video_clip = set_clip_audio(video_clip, audio_clip)

        temp_video = tempfile.NamedTemporaryFile(delete=False, suffix=".mp4")
        temp_video_path = temp_video.name
        temp_video.close()
        
        video_clip.write_videofile(
            temp_video_path,
            fps=24,
            codec="libx264",
            audio_codec="aac",
            verbose=False,
            logger=None
        )
        
        audio_clip.close()
        video_clip.close()

        with open(temp_video_path, "rb") as f:
            mp4_bytes = f.read()

        return mp4_bytes, None

    except Exception as e:
        return None, str(e)

    finally:
        try:
            if temp_audio_path and os.path.exists(temp_audio_path):
                os.remove(temp_audio_path)
            if temp_video_path and os.path.exists(temp_video_path):
                os.remove(temp_video_path)
            for p in temp_img_paths:
                if os.path.exists(p):
                    os.remove(p)
        except Exception:
            pass

def convert_quiz_json_to_markdown(quiz_data: list, lang_name: str) -> str:
    """Converts structured quiz JSON to clean, readable markdown for export."""
    md_out = f"# 🧠 {'Quiz sui Requisiti' if lang_name == 'Italiano' else 'Requirements Quiz'}\n\n"
    for idx, item in enumerate(quiz_data, 1):
        md_out += f"### {idx}. {item.get('question')}\n"
        for opt in item.get('options', []):
            md_out += f"- {opt}\n"
        md_out += "\n"
        
    md_out += f"\n## {'Chiave di Risposta & Spiegazioni' if lang_name == 'Italiano' else 'Answer Key & Explanations'}\n\n"
    for idx, item in enumerate(quiz_data, 1):
        md_out += f"**{idx}. {'Risposta Corretta' if lang_name == 'Italiano' else 'Correct Answer'}:** `{item.get('correct_option')}`\n"
        md_out += f"- **{'Spiegazione' if lang_name == 'Italiano' else 'Explanation'}:** {item.get('explanation')}\n\n"
    return md_out

# ---------------------------------------------------------
# STEP 1: Upload Source Materials (Clean, Plain-Text Indicators)
# ---------------------------------------------------------
ALLOWED_EXTENSIONS = ["docx", "txt", "md", "xlsx", "xls", "csv", "pdf"]

with st.container():
    st.markdown(f"#### {ui['step1_title']}")
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown(f"**{ui['raw_title']}**")
        raw_file = st.file_uploader(
            "Upload FA RAW", 
            type=ALLOWED_EXTENSIONS, 
            key="raw_file_uploader", 
            label_visibility="collapsed",
            help=ui["supported_formats_help"]
        )
        raw_text = parse_uploaded_file(raw_file) if raw_file else ""
        if raw_file:
            st.markdown(
                f"<div class='plain-upload-status'>✅ Loaded: <span class='plain-upload-filename'>{raw_file.name}</span></div>", 
                unsafe_allow_html=True
            )

    with col2:
        st.markdown(f"**{ui['excel_title']}**")
        excel_file = st.file_uploader(
            "Upload Validations", 
            type=ALLOWED_EXTENSIONS, 
            key="excel_file_uploader", 
            label_visibility="collapsed",
            help=ui["supported_formats_help"]
        )
        excel_summary = parse_uploaded_file(excel_file) if excel_file else ""
        if excel_file:
            st.markdown(
                f"<div class='plain-upload-status'>✅ Loaded: <span class='plain-upload-filename'>{excel_file.name}</span></div>", 
                unsafe_allow_html=True
            )

    with col3:
        st.markdown(f"**{ui['figma_title']}**")
        figma_images = st.file_uploader(
            "Upload Screens", 
            type=["png", "jpg", "jpeg"], 
            accept_multiple_files=True, 
            key="figma_uploader", 
            label_visibility="collapsed"
        )
        if figma_images:
            st.markdown(
                f"<div class='plain-upload-status'>✅ Loaded: <span class='plain-upload-filename'>{len(figma_images)} screen(s)</span></div>", 
                unsafe_allow_html=True
            )

st.markdown("---")

# ---------------------------------------------------------
# STEP 2: Deliverable Selection
# ---------------------------------------------------------
st.markdown(f"### {ui['step2_title']}")

if "selected_opt_idx" not in st.session_state:
    st.session_state.selected_opt_idx = 0

btn_col1, btn_col2 = st.columns(2)
for idx, opt_name in enumerate(ui["options"]):
    target_col = btn_col1 if idx < 4 else btn_col2
    with target_col:
        is_active = (st.session_state.selected_opt_idx == idx)
        btn_type = "primary" if is_active else "secondary"
        if st.button(opt_name, key=f"btn_opt_{idx}", type=btn_type):
            st.session_state.selected_opt_idx = idx
            st.rerun()

current_idx = st.session_state.selected_opt_idx
current_option = ui["options"][current_idx]

# ---------------------------------------------------------
# Role Restriction Guardrail
# ---------------------------------------------------------
is_viewer = user_role == ui["roles"][1]
if is_viewer and current_idx < 4:
    st.error(ui["viewer_err"])
    st.stop()

# ---------------------------------------------------------
# Document Matrix Checking
# ---------------------------------------------------------
needs_raw = current_idx in [0, 2, 3, 4, 5, 6]
needs_excel = current_idx in [0, 1, 2, 3, 4, 5, 6]
needs_figma = current_idx in [2, 3, 4]

missing_core_items = []
if needs_raw and not raw_text:
    missing_core_items.append(ui["raw_title"])
if needs_excel and not excel_summary:
    missing_core_items.append(ui["excel_title"])
if needs_figma and not figma_images:
    missing_core_items.append(ui["figma_title"])

# ---------------------------------------------------------
# Centered Orange Circular Progress Slot (Positioned ABOVE the Button)
# ---------------------------------------------------------
progress_slot = st.empty()

# Centered Prominent Action Button
st.markdown("<div style='margin-top: 2rem; margin-bottom: 1.5rem;'></div>", unsafe_allow_html=True)

col_b1, col_b2, col_b3 = st.columns([1, 2, 1])
with col_b2:
    generate_clicked = st.button(
        f"{ui['generate_btn']} : {current_option}",
        key="main_generate_btn",
        type="primary",
        use_container_width=True
    )

if generate_clicked:
    if missing_core_items:
        show_missing_data_popup(missing_core_items)
        st.stop()

    is_quiz_mode = (current_idx == 6)

    system_prompt = f"""
[ROLE & PERSONA]
You are "ReqAssist," an energetic Principal Business Analyst and Product Requirements Architect. Address the user as {user_name}.

[CRUCIAL LANGUAGE RULE]
You MUST respond ENTIRELY and STRICTLY in {lang_key}. All headings, titles, descriptions, table columns, test steps, bullet points, speaker notes, and quiz questions MUST be written in {lang_key}.

[ANTI-HALLUCINATION & STRICT GROUNDING RULES]
1. ZERO INVENTED BUSINESS LOGIC: Do not invent rules not found in inputs.
2. STRICT SOURCE GROUNDING: Ground all answers strictly on the uploaded Functional Notes, Figma UI screens, and Validation sheet.
3. EXPLICIT ASSUMPTION FLAGGING: If an assumption is forced, tag it clearly as '[ASSUMPTION: ...]' (or '[ASSUNZIONE: ...]' if Italian).
4. UNKNOWN DATA: If detail is missing, state: "{'⚠️ Questo dettaglio non è specificato nei documenti forniti. Si prega di consultare il PM/PO.' if lang_key == 'Italiano' else '⚠️ This detail is not specified in the provided documents. Please consult your PM/PO.'}"

[ENHANCED OUTPUT TEMPLATES]

- Acceptance Criteria:
  Generate in TWO distinct, comprehensive sections with HIGH VOLUME of criteria:
  ## {'1. Criteri di Accettazione Funzionali' if lang_key == 'Italiano' else '1. Functional Acceptance Criteria'}
  - Written strictly in BDD format ({'DATO / QUANDO / ALLORA' if lang_key == 'Italiano' else 'GIVEN / WHEN / THEN'}).
  - Must generate an extensive, granular set of criteria covering: Positive Happy Paths, Form Input Validations, Mandatory Fields, Optional Fields, Error Messages, Edge Cases, Empty States, and User Roles/Permissions.
  
  ## {'2. Criteri di Accettazione Tecnici' if lang_key == 'Italiano' else '2. Technical Acceptance Criteria'}
  - MUST ALSO be written strictly in BDD format ({'DATO / QUANDO / ALLORA' if lang_key == 'Italiano' else 'GIVEN / WHEN / THEN'}).
  - Covers API Payload Contracts (200 OK, 400 Bad Request, 401/403 Auth, 500), Database Record Commit & Rollback, Token Header Validation, Response Latency Limits (<500ms), and Structured Error Logging.

- Test Cases:
  Generate in TWO distinct Markdown tables with maximum coverage:
  - Increase the number of test cases significantly. Every single field from the validation sheet must have individual positive, negative, and boundary test cases.
  - Plain text formatting only (no badges/buttons). Include a dedicated Remarks / Note column.
  
  ## {'1. Casi di Test Funzionali' if lang_key == 'Italiano' else '1. Functional Test Cases'}
  | {'ID Test' if lang_key == 'Italiano' else 'Test Case ID'} | {'Campo / Modulo' if lang_key == 'Italiano' else 'Field / Feature Name'} | {'Scenario di Test' if lang_key == 'Italiano' else 'Test Scenario'} | {'Pre-condizioni' if lang_key == 'Italiano' else 'Pre-conditions'} | {'Passaggi di Esecuzione' if lang_key == 'Italiano' else 'Execution Steps'} | {'Risultato Atteso' if lang_key == 'Italiano' else 'Expected Result'} | {'Validazione Campo' if lang_key == 'Italiano' else 'Field Validation Rule'} | {'Tipo (Positivo/Negativo/Edge)' if lang_key == 'Italiano' else 'Test Type (Positive/Negative/Edge)'} | {'Note / Osservazioni' if lang_key == 'Italiano' else 'Remarks / Notes'} |
  
  ## {'2. Casi di Test Tecnici e Integrazione' if lang_key == 'Italiano' else '2. Technical & Integration Test Cases'}
  | {'ID Test Tecnico' if lang_key == 'Italiano' else 'Tech Test ID'} | {'Componente / Endpoint' if lang_key == 'Italiano' else 'Component / Endpoint'} | {'Scenario Tecnico' if lang_key == 'Italiano' else 'Technical Scenario'} | {'Prerequisiti' if lang_key == 'Italiano' else 'Prerequisites'} | {'Payload / Request' if lang_key == 'Italiano' else 'Payload / Request'} | {'Risposta Attesa & Stato DB' if lang_key == 'Italiano' else 'Expected Response & DB State'} | {'Metodo di Verifica' if lang_key == 'Italiano' else 'Verification Method'} | {'Note / Osservazioni' if lang_key == 'Italiano' else 'Remarks / Notes'} |

- Detailed Functional Analysis:
  - Explain the complete, end-to-end user journey in detailed, flowing **PARAGRAPHS ONLY** (avoid sparse bullet lists).
  - Reference the uploaded Figma UI screens, layout sections, step-by-step navigation, state transitions, validation triggers, and system responses in depth.

- Demo Video:
  - Formatted as a complete video production storyboard with a Markdown table:
  | {'Timestamp' if lang_key == 'Italiano' else 'Timestamp'} | {'Schermata / UI Area' if lang_key == 'Italiano' else 'UI Screen / Area'} | {'Azione Visiva' if lang_key == 'Italiano' else 'Visual Action'} | {'Script Voiceover (Voce Narrante)' if lang_key == 'Italiano' else 'Voiceover Script (Narrator)'} |

- PPT Slides (5 to 10 slides):
  - Structured for executive review with Slide Title, Bullets, visual reference note, and full Speaker Notes.

- FAQs:
  - Clean, plain text Markdown format (no button or box structures). Group into Developer Technical FAQs and Stakeholder Business FAQs.

- Quiz Mode:
  - Return a valid JSON array of 8 to 10 question objects with the exact schema:
    [
      {{
        "id": 1,
        "topic": "Short Topic Name (e.g. Validation Rules, API Contract, DB State, RBAC)",
        "question": "Question text in {lang_key}?",
        "options": ["A) ...", "B) ...", "C) ...", "D) ..."],
        "correct_option": "A",
        "explanation": "Detailed explanation grounding why this is correct based on the provided documents."
      }}
    ]
"""

    user_content = f"""
DELIVERABLE REQUESTED: {current_option}
LANGUAGE: {lang_key}

SOURCE MATERIALS:
- FUNCTIONAL NOTES RAW:
{raw_text if raw_text else 'None required or provided'}

- FIELD VALIDATIONS:
{excel_summary if excel_summary else 'None provided'}

Generate the complete artifact strictly adhering to the requested templates and grounding rules.
"""

    contents = [user_content]
    if figma_images and needs_figma:
        for img in figma_images:
            contents.append(Image.open(img))

    try:
        current_pct = 1
        with concurrent.futures.ThreadPoolExecutor() as executor:
            future = executor.submit(
                generate_resilient_content,
                client_inst=client,
                contents=contents,
                system_prompt=system_prompt,
                response_mime_type="application/json" if is_quiz_mode else None
            )

            while not future.done():
                render_circular_progress(progress_slot, current_pct, f"ReqAssist is generating ({lang_key})...")
                if current_pct < 94:
                    current_pct += 1
                time.sleep(0.08)

            output_text, model_used = future.result()

        st.session_state["generated_output"] = output_text
        st.session_state["generated_option"] = current_option
        st.session_state["generated_idx"] = current_idx
        
        st.session_state["generated_mp4_bytes"] = None
        st.session_state["generated_audio_bytes"] = None
        st.session_state["video_error"] = None
        
        # If Demo Video was requested, compile MP4
        if current_idx == 3:
            for p in range(current_pct, 98):
                render_circular_progress(progress_slot, p, "Compiling AI Visual Cards & encoding MP4 video...")
                time.sleep(0.05)
            mp4_bytes, v_err = generate_mp4_video_robust(output_text, figma_images, lang=lang_key)
            if mp4_bytes:
                st.session_state["generated_mp4_bytes"] = mp4_bytes
            else:
                st.session_state["video_error"] = v_err
                audio_bio = extract_voiceover_and_synthesize_audio(output_text, lang=lang_key)
                if audio_bio:
                    st.session_state["generated_audio_bytes"] = audio_bio.getvalue()
        
        if is_quiz_mode:
            try:
                st.session_state["quiz_data"] = json.loads(output_text)
            except Exception:
                st.session_state["quiz_data"] = None
            st.session_state["quiz_submitted"] = False
            st.session_state["quiz_user_answers"] = {}
            
        render_circular_progress(progress_slot, 100, "Finalizing deliverable...")
        time.sleep(0.2)
        progress_slot.empty()
        st.success(f"🎉 Generation Complete! (Powered by {model_used})")

    except Exception as e:
        progress_slot.empty()
        st.error(f"Error during generation: {str(e)}")

# ---------------------------------------------------------
# Output Display & Direct Exports (.docx, .pptx, .md, .xlsx, .mp3, .mp4)
# ---------------------------------------------------------
if "generated_output" in st.session_state and st.session_state["generated_output"]:
    output_text = st.session_state["generated_output"]
    active_opt_name = st.session_state.get("generated_option", current_option)
    active_opt_idx = st.session_state.get("generated_idx", current_idx)

    st.markdown("---")
    tab_out, tab_down = st.tabs([ui["tab_output"], ui["tab_download"]])
    
    with tab_out:
        # 🎬 Demo Video Player
        if active_opt_idx == 3:
            mp4_bytes = st.session_state.get("generated_mp4_bytes", None)
            audio_bytes = st.session_state.get("generated_audio_bytes", None)
            v_err = st.session_state.get("video_error", None)
            
            if mp4_bytes:
                st.markdown(f"#### {ui['video_player_title']}")
                st.video(mp4_bytes, format="video/mp4")
                st.markdown("---")
            elif audio_bytes:
                st.markdown(f"#### {ui['audio_player_title']}")
                st.audio(audio_bytes, format="audio/mp3")
                if v_err:
                    st.info(f"💡 Video renderer notice: {v_err}. Audio voiceover generated as fallback.")
                st.markdown("---")
            elif v_err:
                st.warning(f"⚠️ Video generation notice: {v_err}")
                
        # 🧠 Interactive Playable Quiz Engine
        if active_opt_idx == 6 and st.session_state.get("quiz_data"):
            quiz_items = st.session_state["quiz_data"]
            st.subheader(ui["quiz_header"])
            st.write("")

            if not st.session_state.get("quiz_submitted", False):
                with st.form("interactive_quiz_form"):
                    user_choices = {}
                    for q_idx, item in enumerate(quiz_items):
                        st.markdown(f"##### {q_idx + 1}. {item.get('question')}")
                        opts = item.get("options", [])
                        chosen = st.radio(
                            label=f"Options for Q{q_idx+1}",
                            options=opts,
                            key=f"quiz_radio_{q_idx}",
                            label_visibility="collapsed",
                            index=None
                        )
                        user_choices[q_idx] = chosen
                        st.write("")
                    
                    submitted = st.form_submit_button(ui["quiz_submit_btn"], type="primary", use_container_width=True)
                    if submitted:
                        st.session_state["quiz_user_answers"] = user_choices
                        st.session_state["quiz_submitted"] = True
                        st.rerun()
            else:
                # Calculate Score & Track Missed Areas
                correct_count = 0
                total_count = len(quiz_items)
                missed_topics = []
                
                for q_idx, item in enumerate(quiz_items):
                    user_ans = st.session_state.get("quiz_user_answers", {}).get(q_idx, "")
                    chosen_letter = user_ans.split(")")[0].strip().upper() if user_ans else ""
                    correct_letter = item.get("correct_option", "").strip().upper()
                    if chosen_letter == correct_letter:
                        correct_count += 1
                    else:
                        # Extract topic or question summary
                        topic_name = item.get("topic") or f"Question {q_idx+1}: {item.get('question', '')[:65]}..."
                        missed_topics.append((q_idx + 1, topic_name, item.get("explanation", "")))
                        
                pct = (correct_count / total_count) * 100 if total_count > 0 else 0
                is_passing = (pct >= 75)

                # Emoji selection based on score tier
                if pct >= 90:
                    score_emoji = "🏆"
                elif pct >= 75:
                    score_emoji = "🎯"
                elif pct >= 50:
                    score_emoji = "⚠️"
                else:
                    score_emoji = "❌"

                st.markdown(f"### {ui['quiz_score_title']}")

                # Enhanced Scorecard UI
                theme_color = "#16a34a" if is_passing else "#dc2626"
                bg_color = "#f0fdf4" if is_passing else "#fef2f2"
                border_color = "#86efac" if is_passing else "#fca5a5"

                st.markdown(f"""
                <div style="background-color: {bg_color}; border: 2px solid {border_color}; border-radius: 12px; padding: 22px 28px; margin-bottom: 24px;">
                    <div style="display: flex; align-items: center; justify-content: space-between; flex-wrap: wrap; gap: 16px;">
                        <div>
                            <div style="font-size: 14px; font-weight: 700; color: #475569; text-transform: uppercase; letter-spacing: 0.5px;">Requirement Assessment</div>
                            <div style="font-size: 28px; font-weight: 800; color: #0f172a; margin-top: 4px;">
                                Score: <span style="color: {theme_color};">{correct_count} / {total_count}</span>
                            </div>
                        </div>
                        <div style="text-align: right;">
                            <div style="font-size: 14px; font-weight: 700; color: #475569; text-transform: uppercase; letter-spacing: 0.5px;">Percentage</div>
                            <div style="font-size: 32px; font-weight: 900; color: {theme_color}; margin-top: 2px;">
                                {pct:.1f}% <span style="font-size: 28px;">{score_emoji}</span>
                            </div>
                        </div>
                    </div>
                </div>
                """, unsafe_allow_html=True)

                if is_passing:
                    st.success(ui["quiz_passed"])
                    st.balloons()
                else:
                    st.error(ui["quiz_failed"])
                    
                    # Targeted "Where to focus more" recommendation box
                    st.markdown(f"#### {ui['quiz_focus_title']}")
                    st.markdown(f"*{ui['quiz_focus_intro']}*")
                    for q_num, topic_title, _ in missed_topics:
                        st.markdown(f"- 📌 **Question {q_num} ({topic_title})**")
                    st.write("")

                st.markdown("---")
                st.markdown("#### Detailed Breakdown:")

                for q_idx, item in enumerate(quiz_items):
                    st.markdown(f"##### {q_idx + 1}. {item.get('question')}")
                    opts = item.get("options", [])
                    selected = st.session_state.get("quiz_user_answers", {}).get(q_idx, "")
                    correct_letter = item.get("correct_option", "").strip().upper()
                    
                    chosen_letter = selected.split(")")[0].strip().upper() if selected else ""
                    is_correct = (chosen_letter == correct_letter)
                    
                    for opt in opts:
                        opt_letter = opt.split(")")[0].strip().upper()
                        if opt_letter == correct_letter:
                            st.markdown(f"🟢 **{opt}** *(Correct Answer)*")
                        elif opt_letter == chosen_letter:
                            st.markdown(f"🔴 **{opt}** *(Your Selection)*")
                        else:
                            st.markdown(f"⚪ {opt}")
                            
                    st.info(f"**{ui['quiz_explanation']}** {item.get('explanation')}")
                    st.markdown("---")
                        
                col_r1, col_r2, col_r3 = st.columns([1, 1.5, 1])
                with col_r2:
                    if st.button(ui["quiz_retake_btn"], use_container_width=True):
                        st.session_state["quiz_submitted"] = False
                        st.session_state["quiz_user_answers"] = {}
                        st.rerun()
        else:
            st.markdown(output_text)
        
    with tab_down:
        base_name = active_opt_name.split(" ", 1)[-1].replace(" ", "_").replace("/", "_")
        col_d1, col_d2, col_d3 = st.columns(3)
        
        exportable_text = output_text
        if active_opt_idx == 6 and st.session_state.get("quiz_data"):
            exportable_text = convert_quiz_json_to_markdown(st.session_state["quiz_data"], lang_key)
        
        # 1. Word Document (.docx) Export
        with col_d1:
            docx_bio = create_docx(active_opt_name, exportable_text)
            st.download_button(
                label=ui["docx_btn"],
                data=docx_bio,
                file_name=f"ReqAssist_{base_name}_{user_name}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
            
        # 2. PowerPoint (.pptx) Export (Dark Blue Theme + White Font + Context-Driven Visuals)
        with col_d2:
            try:
                ppt_json_text, _ = generate_resilient_content(
                    client_inst=client,
                    contents=[f"Convert this deliverable into a structured presentation as a JSON array of 6 to 10 slide objects with keys 'title', 'bullets' (list of strings), 'notes' (string). Keep text strictly in {lang_key}:\n\n{exportable_text}"],
                    response_mime_type="application/json"
                )
                slides_data = json.loads(ppt_json_text)
                pptx_bio = create_pptx_with_images(slides_data, figma_files=figma_images if figma_images else None)
                st.download_button(
                    label=ui["pptx_btn"],
                    data=pptx_bio,
                    file_name=f"ReqAssist_{base_name}_{user_name}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
            except Exception:
                st.info("Direct Markdown download available below.")

        # 3. Markdown (.md) Export
        with col_d3:
            st.download_button(
                label=ui["md_btn"],
                data=exportable_text.encode("utf-8"),
                file_name=f"ReqAssist_{base_name}_{user_name}.md",
                mime="text/markdown",
                use_container_width=True
            )
            
        # 4. Excel (.xlsx) Multi-Sheet Export for Test Cases
        if active_opt_idx == 1:
            excel_bio = parse_markdown_tables_to_excel(output_text)
            if excel_bio:
                st.markdown("<br>", unsafe_allow_html=True)
                st.download_button(
                    label=ui["xlsx_btn"],
                    data=excel_bio,
                    file_name=f"ReqAssist_TestCases_{user_name}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True
                )

        # 5. Direct MP4 Video & MP3 Voiceover Audio Downloads for Demo Video
        if active_opt_idx == 3:
            st.markdown("<br>", unsafe_allow_html=True)
            col_v1, col_v2 = st.columns(2)
            
            # MP4 Video Download
            with col_v1:
                mp4_bytes = st.session_state.get("generated_mp4_bytes", None)
                if mp4_bytes:
                    st.download_button(
                        label=ui["mp4_btn"],
                        data=mp4_bytes,
                        file_name=f"ReqAssist_Demo_Walkthrough_{user_name}.mp4",
                        mime="video/mp4",
                        use_container_width=True
                    )
                    
            # MP3 Audio Track Download
            with col_v2:
                audio_bytes = st.session_state.get("generated_audio_bytes", None)
                if not audio_bytes and output_text:
                    audio_bio = extract_voiceover_and_synthesize_audio(output_text, lang=lang_key)
                    audio_bytes = audio_bio.getvalue() if audio_bio else None
                    
                if audio_bytes:
                    st.download_button(
                        label=ui["audio_btn"],
                        data=audio_bytes,
                        file_name=f"ReqAssist_Demo_Voiceover_{user_name}.mp3",
                        mime="audio/mp3",
                        use_container_width=True
                    )
